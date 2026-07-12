"""Turn an uploaded file into indexed evidence.

Agent 1 grounds every requirement in a cited source, so an upload is only useful if we can get
*text* out of it. That means real extraction per format — not storing a blob and pretending.
What we cannot read, we say we cannot read, and we say it at upload time rather than letting the
user discover it three screens later when the agent has nothing to cite.

Audio and video are honestly out of scope here: they need a speech-to-text pass (Gemini or Chirp),
which is a real integration, not a parser. The UI accepts them and marks them
TRANSCRIPTION_PENDING rather than silently dropping them.
"""

from __future__ import annotations

import io
from dataclasses import dataclass

from app.core.logging import log
from app.models import SourceKind

TEXT_EXT = {".txt", ".md", ".csv", ".json", ".eml", ".msg", ".html", ".xml", ".vtt", ".srt"}
MEDIA_EXT = {".mp3", ".wav", ".m4a", ".mp4", ".mov", ".webm", ".ogg"}
IMAGE_EXT = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"}


@dataclass
class Extracted:
    text: str
    kind: SourceKind
    status: str          # EXTRACTED | TRANSCRIPTION_PENDING | OCR_PENDING | UNSUPPORTED
    pages: int = 0
    note: str = ""


def _ext(name: str) -> str:
    return ("." + name.rsplit(".", 1)[-1].lower()) if "." in name else ""


def _guess_kind(name: str, text: str) -> SourceKind:
    n = name.lower()
    if any(k in n for k in ("transcript", "recording", "call", "meeting-audio")):
        return SourceKind.VOICE_TRANSCRIPT
    if any(k in n for k in ("email", "mail", "thread")) or _ext(n) in (".eml", ".msg"):
        return SourceKind.EMAIL
    if any(k in n for k in ("mom", "minutes", "notes", "workshop")):
        return SourceKind.MEETING_NOTES
    return SourceKind.DOCUMENT


def extract(filename: str, data: bytes) -> Extracted:
    ext = _ext(filename)
    try:
        if ext in TEXT_EXT:
            text = data.decode("utf-8", errors="replace")
            return Extracted(text, _guess_kind(filename, text), "EXTRACTED")

        if ext == ".pdf":
            from pypdf import PdfReader

            reader = PdfReader(io.BytesIO(data))
            text = "\n\n".join((p.extract_text() or "") for p in reader.pages).strip()
            if not text:
                # A scanned PDF has no text layer. Reporting "0 requirements found" later would
                # be a lie; saying "this needs OCR" now is the truth.
                return Extracted("", SourceKind.DOCUMENT, "OCR_PENDING", len(reader.pages),
                                 "No text layer — scanned PDF. Needs OCR before it can be cited.")
            return Extracted(text, _guess_kind(filename, text), "EXTRACTED", len(reader.pages))

        if ext in (".docx", ".dotx"):
            from docx import Document

            doc = Document(io.BytesIO(data))
            parts = [p.text for p in doc.paragraphs if p.text.strip()]
            for t in doc.tables:
                for row in t.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        parts.append(" | ".join(cells))
            return Extracted("\n".join(parts), _guess_kind(filename, "\n".join(parts)), "EXTRACTED")

        if ext in (".xlsx", ".xlsm", ".xltx"):
            from openpyxl import load_workbook

            wb = load_workbook(io.BytesIO(data), data_only=True, read_only=True)
            parts = []
            for ws in wb.worksheets:
                parts.append(f"### Sheet: {ws.title}")
                for row in ws.iter_rows(values_only=True):
                    vals = [str(c) for c in row if c is not None]
                    if vals:
                        parts.append(" | ".join(vals))
            return Extracted("\n".join(parts), SourceKind.DOCUMENT, "EXTRACTED")

        if ext in (".pptx", ".potx"):
            from pptx import Presentation

            prs = Presentation(io.BytesIO(data))
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                parts.append(f"### Slide {i}")
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.text_frame.text.strip():
                        parts.append(shape.text_frame.text)
            return Extracted("\n".join(parts), SourceKind.DOCUMENT, "EXTRACTED", len(prs.slides))

        if ext in MEDIA_EXT:
            return Extracted(
                "", SourceKind.VOICE_TRANSCRIPT, "TRANSCRIPTION_PENDING", 0,
                "Audio/video accepted. Speech-to-text is not wired yet — this file cannot be cited "
                "by an agent until it is transcribed.",
            )

        if ext in IMAGE_EXT:
            return Extracted("", SourceKind.DOCUMENT, "OCR_PENDING", 0,
                             "Image accepted. OCR / vision extraction is not wired yet.")

        return Extracted("", SourceKind.DOCUMENT, "UNSUPPORTED", 0,
                         f"No parser for '{ext or 'unknown'}' files.")

    except Exception as e:
        log.exception("ingest.failed", filename=filename)
        return Extracted("", SourceKind.DOCUMENT, "UNSUPPORTED", 0, f"Could not parse: {e}")
